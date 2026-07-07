"""Generate the Prahari FinSpark'26 submission deck (Bank of Maharashtra template).

Run:  python docs/build_deck.py
Out:  docs/Prahari_FinSpark26_Deck.pptx
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---- palette (matches the template) ----
TEAL = RGBColor(0x1C, 0x60, 0x7A)      # titles
INK = RGBColor(0x1E, 0x1E, 0x1E)       # body text
MUTED = RGBColor(0x60, 0x60, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BANNER = RGBColor(0x0E, 0x3C, 0x4A)    # BoM banner box
GOLD = RGBColor(0xF2, 0xC8, 0x3C)
RED = RGBColor(0xD0, 0x3B, 0x3B)
BLUE = RGBColor(0x39, 0x87, 0xE5)
AMBER = RGBColor(0xE0, 0x93, 0x1A)
GREEN = RGBColor(0x0C, 0x8A, 0x3C)
LIGHT = RGBColor(0xEE, 0xF3, 0xF6)
GREY = RGBColor(0xF2, 0xF2, 0xF0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _set(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = font


def textbox(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    _set(tb.text_frame, text, size, color, bold, align)
    return tb


def banner(slide):
    """Bank of Maharashtra tag, top-right (logo placeholder)."""
    b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(10.25), Inches(0.18), Inches(2.85), Inches(0.42))
    b.fill.solid(); b.fill.fore_color.rgb = BANNER; b.line.color.rgb = BANNER
    tf = b.text_frame; tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    _set(tf, "✦  Bank of Maharashtra", 12, GOLD, True, PP_ALIGN.CENTER)


def base(title, subtitle=None):
    s = prs.slides.add_slide(BLANK)
    banner(s)
    textbox(s, Inches(0.55), Inches(0.55), Inches(11.5), Inches(0.9), title, 30, TEAL, True)
    if subtitle:
        textbox(s, Inches(0.6), Inches(1.35), Inches(11.8), Inches(0.9), subtitle, 15, MUTED)
    # thin top rule
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0), Inches(0.08), SW, Inches(0.08))
    ln.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD); ln.line.width = Pt(0.75)
    return s


def bullets(slide, items, x=Inches(0.7), y=Inches(1.7), w=Inches(12.0), h=Inches(5.2),
            size=15, gap=6):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        lvl = 0
        text = it
        if isinstance(it, tuple):
            text, lvl = it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(gap)
        run = p.add_run()
        bullet = "• " if lvl == 0 else "– "
        run.text = bullet + text
        run.font.size = Pt(size if lvl == 0 else size - 1)
        run.font.color.rgb = INK if lvl == 0 else MUTED
        run.font.name = "Calibri"
    return tb


def box(slide, x, y, w, h, text, fill, tcolor=WHITE, size=11, bold=True,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=None):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = line or fill; sp.line.width = Pt(1)
    tf = sp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    _set(tf, text, size, tcolor, bold, PP_ALIGN.CENTER)
    return sp


def arrow(slide, x1, y1, x2, y2, color=TEAL, width=1.75):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color; c.line.width = Pt(width)
    c.line.fill.solid(); c.line.fill.fore_color.rgb = color
    # arrowhead
    le = c.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    tail = le.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    le.append(tail)
    return c


# ============================================================ TITLE
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x0D, 0x0D, 0x0D); bg.line.fill.background()
banner(s)
textbox(s, Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.2), "\U0001F6E1️  PRAHARI", 54, WHITE, True)
textbox(s, Inches(0.85), Inches(3.5), Inches(11.7), Inches(0.8),
        "Privileged-Access Insider-Threat Detection & Management for Banks", 22, GOLD, True)
textbox(s, Inches(0.85), Inches(4.35), Inches(11.7), Inches(0.6),
        "Catch a risky insider — malicious, negligent or compromised — in real time, "
        "respond automatically, and keep the proof quantum-safe.", 14, RGBColor(0xC3, 0xC2, 0xB7))
textbox(s, Inches(0.85), Inches(5.5), Inches(11.7), Inches(0.6),
        "FinSpark'26  ·  Bank of Maharashtra  ·  Problem Statement 1: Privileged Access Misuse & Insider Threat Detection",
        13, WHITE)
textbox(s, Inches(0.85), Inches(6.1), Inches(11.7), Inches(0.5),
        "Team: <your team name>  ·  github.com/positromen/FinSpark-26-Trial", 12, RGBColor(0x89, 0x87, 0x81))

# ============================================================ PROBLEM STATEMENT
s = base("Problem Statement", "PS1 — Privileged Access Misuse & Insider Threat Detection")
bullets(s, [
    "Banks are well defended against outsiders, but the costliest breaches come from trusted insiders — admins, DBAs, IT staff and vendors who already hold privileged access.",
    "The threat has three faces the solution must handle explicitly:",
    ("Malicious — deliberate theft/sabotage (e.g. a dormant vendor bulk-exports customer records at 2 AM).", 1),
    ("Negligent — well-meaning policy breaches (expired vendor access still in use; sensitive data on a personal laptop).", 1),
    ("Compromised — a legitimate account taken over (new geography/device, rapid-fire actions).", 1),
    "Dormant & vendor accounts, privilege escalation, and tamperable audit logs make misuse hard to detect and even harder to prove.",
    "“Harvest-now, decrypt-later”: credentials & audit evidence encrypted with classical crypto today can be broken by future quantum computers.",
    "Why we chose it: highest business impact for a bank (fraud loss, RBI compliance, customer trust) and it lets us fuse real-time AI/UEBA, PAM controls and post-quantum cryptography into one defensible solution.",
])

# ============================================================ SOLUTION OVERVIEW
s = base("Solution Overview", "Prahari (“sentinel”) — what it does")
bullets(s, [
    "A Privileged Access Management + insider-threat platform that sits between privileged staff and critical bank systems and does four things:",
    ("1. WATCHES every privileged action as a normalized, recorded session (replayable command trail).", 1),
    ("2. SCORES each session 0–100 in real time — rule engine + AI behavioural analytics (UEBA) + peer comparison — with a human-readable “why”.", 1),
    ("3. RESPONDS adaptively — allow / step-up MFA / maker-checker / block — tagged with the insider type.", 1),
    ("4. PROTECTS its own credentials & audit log with NIST post-quantum cryptography (ML-KEM-768 vault + ML-DSA-65 signed hash-chain).", 1),
    "Delivered as a real two-sided product: an Employee Portal (staff act, enforcement happens live) and a SOC Console (analysts watch, replay, and review access).",
    "All six judged outcomes delivered: detect misuse ✓  real-time ✓  AI behavioural ✓  risk-based access control ✓  protect admin systems ✓  quantum-proof crypto ✓",
])

# ============================================================ PRE-REQUISITE
s = base("Pre-Requisite", "Assumptions, access, data, tools and inputs")
bullets(s, [
    "Assumption: privileged activity is available as a normalized event stream. For the demo a built-in simulator generates it; in production it is fed from PAM/SIEM.",
    "Required access / data:",
    ("PAM privileged-session logs; IAM user directory (roles, vendor/dormant flags, access-expiry dates); resource/asset inventory.", 1),
    "Per-event inputs: user, action type, resource, records touched, source IP/geo, device fingerprint, timestamp.",
    "Tools / environment:",
    ("Python 3.11+ backend; Node 18+ (front-end build only); a modern browser to view the app.", 1),
    ("Runs fully offline on SQLite — no internet, no external API calls at demo time.", 1),
    ("PQC provider (liboqs) built once on first setup; afterwards entirely offline.", 1),
    "Security prerequisites: `.env`-driven config, no secrets committed, seeded demo data only.",
])

# ============================================================ TOOLS OR RESOURCES
s = base("Tools or Resources", "Technologies, frameworks and libraries used")
bullets(s, [
    "Backend / API: Python 3.14, FastAPI + Uvicorn, SQLAlchemy ORM (SQLite default — Postgres-swappable).",
    "AI / UEBA: scikit-learn IsolationForest (unsupervised anomaly detection) + per-role peer baselining.",
    "Post-quantum crypto: liboqs-python — ML-KEM-768 (FIPS 203) key encapsulation + ML-DSA-65 (FIPS 204) signatures; AES-256-GCM (cryptography) for payloads.",
    "Auth: PBKDF2-HMAC-SHA256 password hashing + HMAC-signed session tokens (standard library — no fragile deps).",
    "Real-time: FastAPI WebSocket push to the SOC console.",
    "Frontend: React 19 + Vite + Tailwind CSS 4 (+ Recharts); dark SOC theme.",
    "Packaging: Docker + docker-compose; one-command run.ps1 / run.sh; prebuilt UI committed for offline demo.",
    "Quality: pytest suite (31 tests — rules, UEBA, scoring, PQC round-trip, auth, live enforcement, scenario distinctness).",
])

# ============================================================ DETECTION & SCORING ENGINE (diagram)
s = base("How It Works — Detection & Scoring Engine")
# pipeline row 1: two detectors -> score -> response
y0 = Inches(1.9)
box(s, Inches(0.7), y0, Inches(3.0), Inches(1.05),
    "RULE ENGINE\nknown-bad patterns\n(weighted, typed)", TEAL, WHITE, 12)
box(s, Inches(0.7), Inches(3.15), Inches(3.0), Inches(1.05),
    "UEBA — IsolationForest\nper-user/role baseline\n+ peer comparison", TEAL, WHITE, 12)
box(s, Inches(4.3), Inches(2.55), Inches(3.0), Inches(1.05),
    "RISK SCORE 0–100\nrules (≤cap) + 0.25×UEBA\n+ peer bonus  → “why”", BLUE, WHITE, 12)
arrow(s, Inches(3.7), Inches(2.42), Inches(4.3), Inches(2.95))
arrow(s, Inches(3.7), Inches(3.67), Inches(4.3), Inches(3.2))
# response ladder
box(s, Inches(8.0), Inches(1.75), Inches(4.6), Inches(0.7), "ALLOW  (score < 40)", GREEN, WHITE, 12)
box(s, Inches(8.0), Inches(2.6), Inches(4.6), Inches(0.7), "STEP-UP MFA  (40–69)", AMBER, RGBColor(0x1A,0x1A,0x19), 12)
box(s, Inches(8.0), Inches(3.45), Inches(4.6), Inches(0.7), "MAKER-CHECKER  (70–84)", RGBColor(0xEC,0x83,0x5A), WHITE, 12)
box(s, Inches(8.0), Inches(4.3), Inches(4.6), Inches(0.7), "BLOCK  (≥ 85)", RED, WHITE, 12)
arrow(s, Inches(7.3), Inches(3.05), Inches(8.0), Inches(3.05))
bullets(s, [
    "Type-aware policy: negligence is floored to maker-checker review and never auto-blocked; malicious blocks; compromised steps up.",
    "Every decision ships a plain-English reason breakdown — explainable AI an auditor can read.",
    "UEBA is a secondary nudge; the stable, explainable rule engine sets the band.",
], y=Inches(5.35), size=13, gap=4)

# ============================================================ THREE INSIDER TYPES (diagram)
s = base("Three Insider Types — One Engine, Three Responses")
lanes = [
    ("MALICIOUS", RED, "Dormant vendor, 2 AM,\nprivilege escalation,\nbulk-export 5000 records",
     "Signals: dormant reactivation,\nprivilege escalation, mass export,\nout-of-role, after-hours",
     "→ BLOCK  (100/100)"),
    ("COMPROMISED", BLUE, "Normal admin account,\nnew country + device,\n6 actions in 90 seconds",
     "Signals: new geo, new device,\natypical hour, rapid-fire\n(inconsistent with human pace)",
     "→ STEP-UP MFA"),
    ("NEGLIGENT", AMBER, "Active vendor, access\nexpired 18 days ago,\nreports from personal laptop",
     "Signals: expired access in use,\nunmanaged/personal device,\nsensitive data touched",
     "→ MAKER-CHECKER review"),
]
x = Inches(0.6)
for name, col, scen, sig, out in lanes:
    w = Inches(3.95)
    box(s, x, Inches(1.85), w, Inches(0.6), name, col, WHITE if col != AMBER else RGBColor(0x1A,0x1A,0x19), 15)
    box(s, x, Inches(2.6), w, Inches(1.15), scen, GREY, INK, 12, bold=False)
    box(s, x, Inches(3.9), w, Inches(1.35), sig, LIGHT, INK, 11, bold=False)
    box(s, x, Inches(5.4), w, Inches(0.7), out, col, WHITE if col != AMBER else RGBColor(0x1A,0x1A,0x19), 13)
    x = Emu(int(x) + int(Inches(4.25)))
textbox(s, Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.8),
        "Same rules + UEBA engine; the dominant rule category sets the insider type, and the response is risk-based on the "
        "nature of the risk — not the score alone.", 13, MUTED)

# ============================================================ POST-QUANTUM SECURITY (diagram)
s = base("Post-Quantum Security Layer", "NIST FIPS 203 / 204 — quantum-safe credentials & evidence")
# Vault flow
textbox(s, Inches(0.7), Inches(1.65), Inches(6), Inches(0.4), "Credential Vault (ML-KEM-768 + AES-256-GCM)", 14, TEAL, True)
box(s, Inches(0.7), Inches(2.15), Inches(2.4), Inches(0.8), "Secret\n(privileged cred)", GREY, INK, 11, False)
box(s, Inches(3.4), Inches(2.15), Inches(2.6), Inches(0.8), "AES-256-GCM\nkey = ML-KEM secret", BLUE, WHITE, 11)
box(s, Inches(0.7), Inches(3.15), Inches(5.3), Inches(0.7), "Stored: ciphertext + KEM ciphertext (decryption needs the KEM secret key)", TEAL, WHITE, 11)
arrow(s, Inches(3.1), Inches(2.55), Inches(3.4), Inches(2.55))
arrow(s, Inches(4.7), Inches(2.95), Inches(4.7), Inches(3.15))
# Audit chain
textbox(s, Inches(6.7), Inches(1.65), Inches(6), Inches(0.4), "Tamper-evident Audit Log (ML-DSA-65)", 14, TEAL, True)
for i, lbl in enumerate(["Entry n-1", "Entry n", "Entry n+1"]):
    bx = Inches(6.7 + i * 2.0)
    box(s, bx, Inches(2.3), Inches(1.75), Inches(1.0),
        f"{lbl}\nhash(prev)\n+ ML-DSA sig", TEAL, WHITE, 10)
    if i < 2:
        arrow(s, Inches(6.7 + i*2.0 + 1.75), Inches(2.8), Inches(6.7 + (i+1)*2.0), Inches(2.8), TEAL, 1.5)
box(s, Inches(6.7), Inches(3.6), Inches(5.75), Inches(0.7),
    "Edit any entry → hash breaks → verify_chain() FAILS at that entry", RED, WHITE, 11)
bullets(s, [
    "store_secret() / get_secret(): AES key is the ML-KEM-768 shared secret — harvest-now-decrypt-later resistant.",
    "append_entry() / verify_chain(): every entry hash-chains the previous and is ML-DSA-65 signed; recomputing hashes isn’t enough — the signature must verify.",
    "Live demo: POST /demo/tamper flips one record → verification fails at the exact entry, on stage.",
], y=Inches(4.6), size=13, gap=4)

# ============================================================ ARCHITECTURE (diagram)
s = base("Architecture Diagram", "Six modules — collection → analytics → scoring → response → PQC")
box(s, Inches(0.5), Inches(1.9), Inches(2.3), Inches(1.0),
    "Privileged activity\n(portal / PAM / SIEM)", MUTED, WHITE, 11)
box(s, Inches(3.1), Inches(1.9), Inches(2.3), Inches(1.0),
    "Collection &\nnormalization\n(Event / Session)", TEAL, WHITE, 11)
box(s, Inches(6.0), Inches(1.35), Inches(2.4), Inches(0.95), "UEBA\nIsolationForest + peer", TEAL, WHITE, 11)
box(s, Inches(6.0), Inches(2.5), Inches(2.4), Inches(0.95), "Rule engine\n(3 insider types)", TEAL, WHITE, 11)
box(s, Inches(9.0), Inches(1.9), Inches(3.2), Inches(1.0), "Risk scoring engine\n0–100 + why + type", BLUE, WHITE, 12)
box(s, Inches(6.0), Inches(4.15), Inches(2.9), Inches(1.0),
    "Adaptive access control\nMFA / maker-checker / block", RGBColor(0xEC,0x83,0x5A), WHITE, 11)
box(s, Inches(9.3), Inches(4.15), Inches(2.9), Inches(1.0),
    "SOC console\nalerts · replay · heatmap · access review", TEAL, WHITE, 11)
box(s, Inches(2.6), Inches(5.5), Inches(8.1), Inches(0.85),
    "Post-Quantum Security Layer  —  ML-KEM-768 credential vault  +  ML-DSA-65 signed, hash-chained audit log", BANNER, GOLD, 12)
arrow(s, Inches(2.8), Inches(2.4), Inches(3.1), Inches(2.4))
arrow(s, Inches(5.4), Inches(2.4), Inches(6.0), Inches(1.85))
arrow(s, Inches(5.4), Inches(2.4), Inches(6.0), Inches(2.97))
arrow(s, Inches(8.4), Inches(1.85), Inches(9.0), Inches(2.3))
arrow(s, Inches(8.4), Inches(2.97), Inches(9.0), Inches(2.5))
arrow(s, Inches(10.6), Inches(2.9), Inches(7.45), Inches(4.15))
arrow(s, Inches(10.6), Inches(2.9), Inches(10.75), Inches(4.15))
textbox(s, Inches(0.5), Inches(6.55), Inches(12.2), Inches(0.6),
        "SQLite by default (zero-friction, offline) — Postgres-swappable via one connection string; the same Event schema ingests real PAM/SIEM feeds.",
        12, MUTED)

# ============================================================ SUPPORTING FUNCTIONAL DOCS
s = base("Supporting Functional Documents", "User flows, process notes and references")
bullets(s, [
    "README.md — quick start, demo accounts, one-command run.",
    "DOCUMENTATION.md — complete technical documentation (architecture, data model, detection math, API, security, deployment) with diagrams.",
    "DEMO_SCRIPT.md — click-by-click 6–7 minute demo narration covering all three insider scenarios, PAM session replay, access review and the audit-tamper moment.",
    "PROJECT_STATUS.md — full project brief, judged-outcome mapping and phase log.",
    "User flows:",
    ("Employee Portal: login → action console → live score → allow / MFA / maker-checker / block.", 1),
    ("SOC Console: login → live sessions → select → why-flagged + timeline + session replay → access review → audit verify/tamper.", 1),
    "Logic flows (diagrams in docs): detection→score→response, ML-KEM vault, ML-DSA hash-chain, three-scenario sequences.",
])

# ============================================================ KEY DIFFERENTIATORS & ADOPTION
s = base("Key Differentiators & Adoption Plan")
bullets(s, [
    "Differentiators:",
    ("All THREE insider types (malicious / negligent / compromised) with distinct, correct responses — most tools do anomaly detection only.", 1),
    ("A real PAM surface: privileged-session recording (replayable command trail) + access review (dormant / vendor / expired) — not just alerts.", 1),
    ("Explainable, graduated enforcement applied LIVE (blocks the action mid-flight), not a passive dashboard.", 1),
    ("Working NIST post-quantum crypto for credentials and a tamper-evident audit log.", 1),
    "Adoption plan (low-risk rollout):",
    ("Phase 1 — deploy read-only in monitor mode beside existing PAM/SIEM; learn baselines, no enforcement.", 1),
    ("Phase 2 — enable step-up MFA + maker-checker for medium risk; SOC runs replay & access reviews.", 1),
    ("Phase 3 — auto-block high-confidence malicious; move crown-jewel credentials into the PQC vault.", 1),
    ("Integrations — CyberArk / BeyondTrust PAM, Splunk / QRadar SIEM, bank IdP / SSO.", 1),
])

# ============================================================ UNIQUENESS
s = base("Uniqueness of Approach and Solution")
bullets(s, [
    "One engine, three insider types — and a type-aware policy: negligence goes to human review, never an automated block. Response is risk-based on the nature of the risk, not the score alone.",
    "Rules + UEBA fused into a single explainable 0–100 score with human-readable “why” — no black box.",
    "PAM + UEBA + post-quantum crypto combined in one offline-capable product — a combination competitors don’t ship together.",
    "Live enforcement via a real employee portal that stops the action in progress and locks the account — challenged actions aren’t even persisted until allowed.",
    "Post-quantum, tamper-evident audit chain — you literally cannot silently edit the log; forging requires breaking ML-DSA-65.",
    "Session recording replays the exact privileged command trail, with blocked commands struck through — forensic-grade evidence.",
])

# ============================================================ USER EXPERIENCE
s = base("User Experience", "Two role-based experiences behind one login")
bullets(s, [
    "Employee Portal — a clean action console: pick a system, choose an action, act. Risk feedback is immediate and self-explanatory:",
    ("Amber step-up MFA modal · orange “held for approval” · full-screen red BLOCK overlay with reasons.", 1),
    "SOC Console — single-screen situational awareness: live sessions, risk gauge, insider-type-tagged alerts, why-flagged panel, session timeline, terminal-style session replay, users×day risk heatmap, PAM access review.",
    "Color-coded insider types (red / blue / amber); new alerts flash; every flag carries a plain-English reason for fast triage.",
    "Professional dark SOC theme; tabular numerics; responsive; no 3D, no clutter — optimized for a control-room wall.",
    "Accessibility: status is never colour-alone (icon + label), high-contrast surfaces, keyboard-friendly controls.",
])

# ============================================================ BUSINESS POTENTIAL
s = base("Business Potential and Relevance")
bullets(s, [
    "Insider-driven incidents are among the costliest breach categories — one misused privileged account can expose millions of customer records and trigger regulatory penalties.",
    "Directly supports RBI cyber-security & governance expectations: least-privilege, immutable audit trails, privileged-access oversight.",
    "Quantum-readiness aligns with global PQC migration mandates (NIST FIPS 203/204) — protects long-lived secrets and evidence today.",
    "Expansion: PS2 transaction-fraud correlation, periodic entitlement reviews, just-in-time access, insider-risk scoring for the whole workforce.",
    "Market: every bank, NBFC and fintech with privileged users; deployable as an internal platform or SOC-as-a-service.",
    "Measurable value: faster mean-time-to-detect, fewer standing privileges, provable audit integrity, reduced breach exposure.",
])

# ============================================================ SCALABILITY
s = base("Scalability", "From a laptop demo to a bank-wide deployment")
bullets(s, [
    "Stateless scoring service — scale horizontally behind a load balancer; each request scores one session.",
    "Storage: SQLite → PostgreSQL by changing one connection string (pure ORM, no raw SQL).",
    "Real-time fan-out: WebSocket behind a Redis / Kafka broker to serve many SOC clients and high event throughput.",
    "Ingestion generalizes to enterprise PAM/SIEM feeds (millions of privileged events/day) mapping onto the same Event schema.",
    "ML: models retrain per role / shift nightly; a feature store holds baselines; the UEBA interface accepts an autoencoder upgrade with no app changes.",
    "PQC operations (ML-KEM encapsulate, ML-DSA sign) are sub-millisecond — negligible next to a DB write.",
])

# ============================================================ EASE OF DEPLOYMENT
s = base("Ease of Deployment and Maintenance")
bullets(s, [
    "One command — run.ps1 (Windows) or run.sh (Linux/macOS) — creates the venv, installs deps, seeds the DB, builds/serves the UI and API, fully offline.",
    "Container-ready: Docker + docker-compose for repeatable deploys.",
    "Config via .env / pydantic settings; no secrets in the repo; safe defaults.",
    "Small, modular, type-hinted codebase — one concern per module (rules, ueba, score, response, live, pam, security/*).",
    "31 automated tests (pytest) guard detection logic, scoring bands, PQC round-trip, auth and live enforcement — safe to iterate.",
    "PQC behind a single abstraction (app/security/pqc.py) — swap providers without touching application code.",
    "Prebuilt front-end committed — the demo laptop needs no Node/npm at the venue.",
])

# ============================================================ SECURITY CONSIDERATIONS
s = base("Security Considerations")
bullets(s, [
    "Data minimization: scores on privileged-activity metadata — no customer PII required.",
    "Identity: PBKDF2-HMAC-SHA256 password hashing; HMAC-signed session tokens; role-gated APIs (employee vs analyst; 403 on cross-role).",
    "Credential protection: ML-KEM-768-wrapped AES-256-GCM vault for privileged secrets (quantum-safe).",
    "Audit integrity: hash-chained + ML-DSA-65-signed entries; tampering is cryptographically evident (proven live).",
    "Enforcement is server-side: blocked accounts stay locked; a challenged action is not persisted until actually allowed (no score-gaming, no bypass by re-login).",
    "Offline by design — no external calls at runtime; keys kept out of git; least-privilege demo accounts.",
    "Compliance-friendly: explainable decisions, immutable trail, access reviews for standing/dormant/expired grants.",
])

# ============================================================ GITHUB REPO LINK
s = base("GitHub Repository & Supporting Material")
bullets(s, [
    "Repository:  github.com/positromen/FinSpark-26-Trial",
    "One-command offline run (run.ps1 / run.sh); prebuilt UI committed; 31 passing tests.",
    "In the repo: full source, DOCUMENTATION.md (with diagrams), DEMO_SCRIPT.md, this deck, Dockerfile + docker-compose.",
    "Demo accounts (password prahari123): soc_admin (SOC), rmehta (DBA), ext_dsouza (dormant vendor / attacker), ext_rao (expired vendor / negligent). MFA code 246810.",
], size=15)
box(s, Inches(0.7), Inches(4.4), Inches(11.9), Inches(1.7),
    "Architecture: Privileged activity → Collection → (UEBA + Rules) → Risk score 0–100 → Adaptive response + SOC console"
    "  — all wrapped by the ML-KEM / ML-DSA post-quantum security layer.", LIGHT, INK, 13, False)

# ============================================================ SCREENSHOTS & VIDEO
s = base("Solution Screenshots, Video & GitHub Link", "Proof the solution is functional")
bullets(s, [
    "Screenshots to include (placeholders — drop images in here):",
    ("Login · Employee Portal action console · full-screen BLOCK overlay.", 1),
    ("SOC Console: live blocked session + insider-type-tagged CRITICAL alert + why-flagged panel.", 1),
    ("Privileged-session recording replay (DENIED command struck through).", 1),
    ("PAM access review (dormant / vendor / expired flags) · risk heatmap.", 1),
    ("Audit chain VERIFIED (green) → after Tamper: FAILED (red).", 1),
    "Demo video:  <add screen-recording link>",
    "GitHub:  github.com/positromen/FinSpark-26-Trial",
])
for i in range(4):
    ph = box(s, Inches(0.7 + i*3.05), Inches(5.55), Inches(2.8), Inches(1.4),
             "screenshot\nplaceholder", GREY, MUTED, 11, False)

prs.save("docs/Prahari_FinSpark26_Deck.pptx")
print(f"Saved docs/Prahari_FinSpark26_Deck.pptx with {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
